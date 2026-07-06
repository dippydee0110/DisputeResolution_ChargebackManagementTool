from __future__ import annotations

import re
from datetime import datetime
from pathlib import Path
from typing import List, Tuple

from reportlab.lib import colors
from reportlab.lib.pagesizes import LETTER
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.utils import ImageReader
from reportlab.platypus import Image, PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


PRIMARY = colors.HexColor("#0A2540")
ACCENT = colors.HexColor("#1177CC")
MUTED = colors.HexColor("#5F6B7A")
LIGHT_BG = colors.HexColor("#F4F7FB")


def sanitize_text(text: str) -> str:
    # Remove interview-related wording from generated artifacts.
    cleaned = re.sub(r"\binterview[- ]?ready\b", "program-ready", text, flags=re.IGNORECASE)
    cleaned = re.sub(r"\binterview\b", "program", cleaned, flags=re.IGNORECASE)
    return cleaned


def get_visual_assets(output_dir: Path) -> List[Tuple[Path, str]]:
    folder = output_dir / "visual_evidence"
    visuals = [
        (folder / "01_agentic_workflow.png", "Agentic Workflow Diagram: early-dispute detection and classification flow"),
        (folder / "02_agent_loop_workflow.png", "Agent Loop Workflow: LLM planning, tools, and policy gate controls"),
        (folder / "03_merchant_risk_dashboard.png", "Merchant Risk Dashboard: operations queue, evidence intake, and actions"),
        (folder / "04_customer_expanded_case.png", "Customer Expanded Case: AI summary, money movement, and lifecycle reference"),
        (folder / "05_merchant_expanded_case.png", "Merchant Expanded Case: recommendation-led investigation and evidence trace"),
        (folder / "06_admin_expanded_case.png", "Admin Expanded Case: intervention queue, settlement view, and lifecycle controls"),
    ]
    return [(path, caption) for path, caption in visuals if path.exists()]


def parse_markdown_lines(md_text: str) -> List[Tuple[str, str]]:
    items: List[Tuple[str, str]] = []
    for raw in sanitize_text(md_text).splitlines():
        line = raw.rstrip()
        if not line.strip():
            items.append(("blank", ""))
            continue
        if line.startswith("### "):
            items.append(("h3", line[4:].strip()))
            continue
        if line.startswith("## "):
            items.append(("h2", line[3:].strip()))
            continue
        if line.startswith("# "):
            items.append(("h1", line[2:].strip()))
            continue
        if line.startswith("- "):
            items.append(("bullet", line[2:].strip()))
            continue
        if line[:2].isdigit() and line[2:4] == ". ":
            items.append(("num", line[4:].strip()))
            continue
        items.append(("p", line.strip()))
    return items


def build_pdf(md_path: Path, pdf_path: Path, doc_title: str) -> None:
    styles = getSampleStyleSheet()
    style_h1 = ParagraphStyle("H1X", parent=styles["Heading1"], fontName="Helvetica-Bold", textColor=PRIMARY, fontSize=20, leading=24, spaceAfter=8)
    style_h2 = ParagraphStyle("H2X", parent=styles["Heading2"], fontName="Helvetica-Bold", textColor=PRIMARY, fontSize=14, leading=18, spaceBefore=10, spaceAfter=5)
    style_h3 = ParagraphStyle("H3X", parent=styles["Heading3"], fontName="Helvetica-Bold", textColor=ACCENT, fontSize=11, leading=14, spaceBefore=8, spaceAfter=3)
    style_p = ParagraphStyle("PX", parent=styles["BodyText"], fontName="Helvetica", textColor=colors.black, fontSize=10, leading=14, spaceAfter=5)
    style_bullet = ParagraphStyle("BX", parent=style_p, leftIndent=14, bulletIndent=4)
    style_meta = ParagraphStyle("META", parent=style_p, fontSize=9, textColor=MUTED)
    style_caption = ParagraphStyle("CAPTION", parent=style_p, fontSize=9, textColor=MUTED, spaceBefore=2, spaceAfter=6)

    content = md_path.read_text(encoding="utf-8")
    parsed = parse_markdown_lines(content)
    visuals = get_visual_assets(md_path.parent)

    story = [
        Paragraph(doc_title, style_h1),
        Paragraph("Disputes & Chargebacks Program Artifact", style_meta),
        Paragraph(datetime.now().strftime("Generated %Y-%m-%d %H:%M"), style_meta),
        Spacer(1, 8),
    ]

    highlights = [
        ["P0", "Implement canonical lifecycle, policy guardrails, and stage-level evidence."],
        ["P1", "Scale to modular services for decisioning, compliance, and ledger integrity."],
        ["Metrics", "Improve win rates, reduce net loss, and preserve host/guest trust KPIs."],
    ]
    highlight_table = Table(highlights, colWidths=[50, 460])
    highlight_table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, -1), LIGHT_BG),
                ("TEXTCOLOR", (0, 0), (0, -1), PRIMARY),
                ("TEXTCOLOR", (1, 0), (1, -1), colors.black),
                ("FONTNAME", (0, 0), (0, -1), "Helvetica-Bold"),
                ("FONTNAME", (1, 0), (1, -1), "Helvetica"),
                ("FONTSIZE", (0, 0), (-1, -1), 9),
                ("GRID", (0, 0), (-1, -1), 0.3, colors.HexColor("#D7E1EC")),
                ("LEFTPADDING", (0, 0), (-1, -1), 6),
                ("RIGHTPADDING", (0, 0), (-1, -1), 6),
                ("TOPPADDING", (0, 0), (-1, -1), 4),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
            ]
        )
    )
    story.extend([highlight_table, Spacer(1, 10)])

    for kind, text in parsed:
        if kind == "blank":
            story.append(Spacer(1, 4))
        elif kind == "h1":
            story.append(Paragraph(text, style_h1))
        elif kind == "h2":
            story.append(Paragraph(text, style_h2))
        elif kind == "h3":
            story.append(Paragraph(text, style_h3))
        elif kind in ("bullet", "num"):
            story.append(Paragraph(f"• {text}", style_bullet))
        else:
            story.append(Paragraph(text, style_p))

    if visuals:
        story.extend([PageBreak(), Paragraph("Visual Evidence Appendix", style_h2), Spacer(1, 4)])
        max_width = 540
        max_height = 300
        for image_path, caption in visuals:
            img_reader = ImageReader(str(image_path))
            img_w, img_h = img_reader.getSize()
            scale = min(max_width / img_w, max_height / img_h)
            width = img_w * scale
            height = img_h * scale
            story.append(Paragraph(caption, style_h3))
            story.append(Image(str(image_path), width=width, height=height))
            story.append(Paragraph(f"File: {image_path.name}", style_caption))
            story.append(Spacer(1, 6))

    footer_table = Table([["Generated for interview package", "Source: " + md_path.name]], colWidths=[250, 260])
    footer_table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, -1), LIGHT_BG),
                ("TEXTCOLOR", (0, 0), (-1, -1), MUTED),
                ("FONTSIZE", (0, 0), (-1, -1), 8),
                ("LEFTPADDING", (0, 0), (-1, -1), 6),
                ("RIGHTPADDING", (0, 0), (-1, -1), 6),
                ("TOPPADDING", (0, 0), (-1, -1), 4),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
            ]
        )
    )
    footer_table._cellvalues[0][0] = "Generated for disputes program artifacts"
    story.append(Spacer(1, 8))
    story.append(footer_table)

    doc = SimpleDocTemplate(str(pdf_path), pagesize=LETTER, leftMargin=36, rightMargin=36, topMargin=36, bottomMargin=36)
    doc.build(story)


def chunk_content_for_slides(parsed: List[Tuple[str, str]], max_lines: int = 11) -> List[List[str]]:
    chunks: List[List[str]] = []
    current: List[str] = []

    def flush() -> None:
        nonlocal current
        if current:
            chunks.append(current)
            current = []

    for kind, text in parsed:
        if kind == "blank":
            continue
        if kind in ("h1", "h2", "h3"):
            line = text.upper() if kind == "h1" else text
        elif kind in ("bullet", "num"):
            line = f"• {text}"
        else:
            line = text

        if len(current) >= max_lines:
            flush()
        current.append(line)

    flush()
    return chunks


def build_presentation(sources: List[Tuple[Path, str]], pptx_path: Path) -> None:
    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Chargebacks and Disputes Program Plan"
    title_slide.placeholders[1].text = "Professional program package: P0 execution, P1 architecture, and metrics targets"

    title_shape = title_slide.shapes.title
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(10, 37, 64)
    subtitle = title_slide.placeholders[1].text_frame.paragraphs[0]
    subtitle.font.size = Pt(16)
    subtitle.font.color.rgb = RGBColor(47, 84, 117)

    def style_content_slide(slide, heading: str) -> None:
        title = slide.shapes.title
        title.text = heading
        p = title.text_frame.paragraphs[0]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(10, 37, 64)

    for md_path, section_title in sources:
        divider = prs.slides.add_slide(prs.slide_layouts[2])
        style_content_slide(divider, section_title)
        divider.placeholders[1].text = f"Source: {md_path.name}"
        divider.placeholders[1].text_frame.paragraphs[0].font.size = Pt(14)
        divider.placeholders[1].text_frame.paragraphs[0].font.color.rgb = RGBColor(60, 73, 90)

        parsed = parse_markdown_lines(md_path.read_text(encoding="utf-8"))
        chunks = chunk_content_for_slides(parsed)

        for idx, chunk in enumerate(chunks, start=1):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            style_content_slide(slide, f"{section_title} ({idx}/{len(chunks)})")
            tf = slide.placeholders[1].text_frame
            tf.clear()

            for line_index, line in enumerate(chunk):
                if line_index == 0:
                    tf.text = line
                    tf.paragraphs[0].font.size = Pt(18)
                    tf.paragraphs[0].font.color.rgb = RGBColor(21, 45, 74)
                else:
                    p = tf.add_paragraph()
                    p.text = line
                    p.font.size = Pt(16)
                    p.font.color.rgb = RGBColor(44, 62, 80)

    visuals = get_visual_assets(pptx_path.parent)
    if visuals:
        section = prs.slides.add_slide(prs.slide_layouts[2])
        style_content_slide(section, "Visual Evidence")
        section.placeholders[1].text = "Operational evidence from customer, merchant, admin, and workflow views"
        section.placeholders[1].text_frame.paragraphs[0].font.size = Pt(14)
        section.placeholders[1].text_frame.paragraphs[0].font.color.rgb = RGBColor(60, 73, 90)

        for image_path, caption in visuals:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12), Inches(0.6))
            p = tx.text_frame.paragraphs[0]
            p.text = caption
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = RGBColor(10, 37, 64)
            p.alignment = PP_ALIGN.LEFT

            slide.shapes.add_picture(str(image_path), Inches(0.7), Inches(1.1), width=Inches(12.0), height=Inches(5.8))

    prs.save(str(pptx_path))


def main() -> None:
    root = Path(__file__).resolve().parents[1]
    output = root / "output"

    p0 = output / "2026-07-05-P0-implementation-backlog.md"
    p1 = output / "2026-07-05-P1-architecture-roadmap.md"
    metrics = output / "2026-07-05-metrics-results-slide.md"

    build_pdf(p0, output / "2026-07-05-P0-implementation-backlog.pdf", "P0 Implementation Backlog")
    build_pdf(p1, output / "2026-07-05-P1-architecture-roadmap.pdf", "P1 Architecture Roadmap")
    build_pdf(metrics, output / "2026-07-05-metrics-results-slide.pdf", "Metrics and Results")

    build_presentation(
        [
            (p0, "P0 Implementation Backlog"),
            (p1, "P1 Architecture Roadmap"),
            (metrics, "Metrics and Results"),
        ],
        output / "2026-07-05-disputes-program-plan.pptx",
    )

    print("Created:")
    print(output / "2026-07-05-P0-implementation-backlog.pdf")
    print(output / "2026-07-05-P1-architecture-roadmap.pdf")
    print(output / "2026-07-05-metrics-results-slide.pdf")
    print(output / "2026-07-05-disputes-program-plan.pptx")


if __name__ == "__main__":
    main()
