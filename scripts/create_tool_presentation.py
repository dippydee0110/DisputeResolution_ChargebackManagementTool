from pathlib import Path
import struct
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

OUT_PATH = Path(r"c:\Users\suman\Documents\Dispute & Chargeback Management Agent\output\TOOL_OVERVIEW_PRESENTATION.pptx")
OUT_PATH.parent.mkdir(parents=True, exist_ok=True)

SCREENSHOT_DIR = Path(r"c:\Users\suman\Documents\Dispute & Chargeback Management Agent\output\screenshots")
OUTPUT_DIR = Path(r"c:\Users\suman\Documents\Dispute & Chargeback Management Agent\output")

prs = Presentation()

THEME_PRIMARY = RGBColor(34, 49, 63)
THEME_ACCENT = RGBColor(31, 111, 235)
THEME_TEXT = RGBColor(34, 34, 34)
THEME_WHITE = RGBColor(255, 255, 255)


def get_png_dimensions(image_path: Path):
    with image_path.open("rb") as f:
        header = f.read(24)
    if len(header) < 24 or header[0:8] != b"\x89PNG\r\n\x1a\n":
        return None
    width, height = struct.unpack(">II", header[16:24])
    if width <= 0 or height <= 0:
        return None
    return width, height


def add_fitted_picture(slide, image_path: Path, left: float, top: float, max_width: float, max_height: float):
    dims = get_png_dimensions(image_path)
    if not dims:
        slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), Inches(max_width), Inches(max_height))
        return

    img_w, img_h = dims
    img_ratio = img_w / img_h
    box_ratio = max_width / max_height

    if img_ratio >= box_ratio:
        draw_w = max_width
        draw_h = max_width / img_ratio
    else:
        draw_h = max_height
        draw_w = max_height * img_ratio

    draw_left = left + (max_width - draw_w) / 2
    draw_top = top + (max_height - draw_h) / 2
    slide.shapes.add_picture(str(image_path), Inches(draw_left), Inches(draw_top), Inches(draw_w), Inches(draw_h))


def apply_theme(slide):
    top_band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.33),
        Inches(0.6),
    )
    top_band.fill.solid()
    top_band.fill.fore_color.rgb = THEME_PRIMARY
    top_band.line.fill.background()


def style_title(shape):
    tf = shape.text_frame
    for p in tf.paragraphs:
        if not p.runs:
            continue
        for run in p.runs:
            run.font.color.rgb = THEME_TEXT
            run.font.bold = True
            run.font.size = Pt(34)


def add_notes(slide, notes):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes


def add_title_slide(title, subtitle, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    apply_theme(slide)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    style_title(slide.shapes.title)
    for p in slide.placeholders[1].text_frame.paragraphs:
        for r in p.runs:
            r.font.color.rgb = THEME_ACCENT
            r.font.size = Pt(18)
    add_notes(slide, notes)


def add_bullets_slide(title, bullets, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_theme(slide)
    slide.shapes.title.text = title
    style_title(slide.shapes.title)
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = THEME_TEXT
    add_notes(slide, notes)


def add_two_column_slide(title, left_title, left_points, right_title, right_points, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    apply_theme(slide)
    slide.shapes.title.text = title
    style_title(slide.shapes.title)

    left_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.2))
    left_tf = left_box.text_frame
    left_tf.text = left_title
    left_tf.paragraphs[0].font.bold = True
    left_tf.paragraphs[0].font.size = Pt(24)
    left_tf.paragraphs[0].font.color.rgb = THEME_ACCENT
    for point in left_points:
        p = left_tf.add_paragraph()
        p.text = point
        p.level = 1
        p.font.size = Pt(20)
        p.font.color.rgb = THEME_TEXT

    right_box = slide.shapes.add_textbox(Inches(6.6), Inches(1.5), Inches(5.8), Inches(5.2))
    right_tf = right_box.text_frame
    right_tf.text = right_title
    right_tf.paragraphs[0].font.bold = True
    right_tf.paragraphs[0].font.size = Pt(24)
    right_tf.paragraphs[0].font.color.rgb = THEME_ACCENT
    for point in right_points:
        p = right_tf.add_paragraph()
        p.text = point
        p.level = 1
        p.font.size = Pt(20)
        p.font.color.rgb = THEME_TEXT

    add_notes(slide, notes)


def add_image_slide(title, image_name, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    apply_theme(slide)
    slide.shapes.title.text = title
    style_title(slide.shapes.title)

    image_path = SCREENSHOT_DIR / image_name
    if image_path.exists():
        add_fitted_picture(slide, image_path, 0.8, 1.3, 11.8, 5.8)
    else:
        fallback = slide.shapes.add_textbox(Inches(1.2), Inches(2.5), Inches(10.8), Inches(2.0))
        tf = fallback.text_frame
        tf.text = f"Screenshot not found: {image_name}"
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.color.rgb = THEME_TEXT

    add_notes(slide, notes)


def add_output_image_slide(title, image_name, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    apply_theme(slide)
    slide.shapes.title.text = title
    style_title(slide.shapes.title)

    image_path = OUTPUT_DIR / image_name
    if image_path.exists():
        add_fitted_picture(slide, image_path, 0.8, 1.3, 11.8, 5.8)
    else:
        fallback = slide.shapes.add_textbox(Inches(1.2), Inches(2.5), Inches(10.8), Inches(2.0))
        tf = fallback.text_frame
        tf.text = f"Diagram not found: {image_name}"
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.color.rgb = THEME_TEXT

    add_notes(slide, notes)


add_title_slide(
    "Early Dispute Detection and Chargeback Management Tool",
    "Presentation Overview | Generated 2026-07-01",
    "Introduce this as an end-to-end platform for pre-chargeback detection, resolution workflow, and administrative control.",
)

add_bullets_slide(
    "What This Tool Does",
    [
        "Detects dispute signals before formal chargeback filing",
        "Scores case risk to prioritize operational action",
        "Supports customer and merchant AI-assisted live chat negotiation",
        "Separates disputes from official chargebacks",
        "Enables evidence prompts and automatic closure when both parties agree",
    ],
    "Explain that early detection plus workflow visibility is the core value proposition.",
)

add_bullets_slide(
    "Core System Components",
    [
        "Data generation layer: structured chargeback records with reason-code fields",
        "Detection layer: chat and evidence signal extraction with weighted risk scoring",
        "Resolution layer: role-specific customer and merchant workflows",
        "Governance layer: admin intervention, policy checks, and escalation controls",
        "Reporting layer: operational dashboards and integration-ready output files",
    ],
    "Present this as a layered operating model so stakeholders understand ownership by function.",
)

add_bullets_slide(
    "Function Overview (Scripts)",
    [
        "generate_chargeback_dataset.py: creates normalized sample case data for reporting and QA",
        "early_dispute_detection.py: detects early dispute indicators and computes risk classifications",
        "run_analysis.py: orchestrates full execution and regenerates all downstream deliverables",
        "create_tool_overview_pdf.py: produces stakeholder-ready summary PDF with screenshots",
        "create_tool_presentation.py: builds this presentation with theme, visuals, and speaker notes",
    ],
    "Keep this slide concise and implementation-focused: what each script does and why it exists.",
)

add_two_column_slide(
    "Dispute vs Chargeback Lifecycle",
    "Dispute Stage",
    [
        "Customer or merchant can initiate",
        "Negotiation with notes, evidence, and amount updates",
        "Accept or reject outcomes tracked",
        "Risk level drives response priority",
    ],
    "Chargeback Stage",
    [
        "Officially filed at issuer/network level",
        "Reason code validation and evidence arbitration",
        "Admin escalation and resolution outcomes",
        "Final result tracked as won or lost",
    ],
    "Clarify the lifecycle boundary and why this separation reduces confusion and delays.",
)

add_bullets_slide(
    "Risk and Detection Logic",
    [
        "Risk tiers: Critical, High, Medium, Low",
        "Signal weighting from conversation patterns",
        "Examples: unrecognized transaction, refund demand, non-delivery",
        "Reason-code aware triage for faster routing",
        "High-risk alerts support proactive outreach",
    ],
    "Mention that thresholds are configurable and should be tuned by merchant segment.",
)

add_bullets_slide(
    "Dashboards and User Views",
    [
        "Customer Resolution Center: active chat, evidence prompts, and settlement response",
        "Merchant Resolution Center: live evidence exchange and representment readiness",
        "Merchant Admin Dashboard: intervention queue, risk segmentation, and policy oversight",
        "Disputes and chargebacks are separated to preserve lifecycle clarity",
        "Visual analytics track risk distribution, case aging, and exposure by status",
    ],
    "Explain how each interface supports a distinct operational role with clear handoffs.",
)

add_image_slide(
    "Application Context: Entry and Navigation",
    "index.png",
    "Use this view to explain the information architecture: role-based entry points and guided navigation to each workflow.",
)

add_image_slide(
    "Screenshot: Merchant Admin Dashboard",
    "merchant_dashboard.png",
    "Highlight merchant risk drilldown, workflow controls, evidence readiness actions, and issuer outcome simulation.",
)

add_image_slide(
    "Screenshot: Early Dispute Detection",
    "early_dispute_detection.png",
    "Highlight risk-ranked cases, signal traceability, and triage-ready recommendations for frontline teams.",
)

add_output_image_slide(
    "Workflow Diagram: Agentic Detection Flow",
    "agentic_workflow_diagram.png",
    "Walk through the end-to-end logic: chat input, signal detection, weighted scoring, risk classification, and output generation.",
)

add_output_image_slide(
    "Workflow Diagram: Agent Loop (LLM + Tools)",
    "agent_loop_llm_tools_workflow.png",
    "Explain the agent loop: LLM planning, tool execution, confidence/policy gating, human escalation, action execution, and closed-loop state persistence.",
)

add_bullets_slide(
    "Workflow Interpretation",
    [
        "Input context is captured from customer-support conversations and evidence artifacts",
        "Signals are matched once per category to reduce duplicate inflation",
        "Weighted scoring determines priority bands for operational action",
        "Risk levels map directly to response urgency and ownership",
        "AI prompt chips guide evidence collection and confirmation in live chat",
        "Cases auto-close when both parties explicitly agree in the chat flow",
        "Outputs support both daily operations and system integration",
    ],
    "Frame this as a governance-friendly decision pipeline: transparent, explainable, and actionable.",
)

add_image_slide(
    "Screenshot: Resolution Center (Customer)",
    "resolution_center_customer.png",
    "Show customer-side live chat execution with evidence prompts, amount context, and agreement-based closure controls.",
)

add_image_slide(
    "Screenshot: Resolution Center (Merchant)",
    "resolution_center_merchant.png",
    "Show merchant-side live chat response controls, evidence exchange, and mutual-agreement closure progression.",
)

add_bullets_slide(
    "Outputs and Deliverables",
    [
        "Excel dataset for case operations and audits",
        "Interactive HTML dashboards for daily workflows",
        "JSON output for downstream system integration",
        "Markdown summary for run-level insights",
        "Export-ready reporting for stakeholders",
    ],
    "Connect each output to its consumer: analyst, operator, or system integration.",
)

add_bullets_slide(
    "Admin Intervention Workflow",
    [
        "Request evidence",
        "Escalate to network review",
        "Mark chargeback won",
        "Mark chargeback lost",
        "Track status and rationale for compliance",
    ],
    "Stress auditability and policy-driven intervention consistency.",
)

add_bullets_slide(
    "Business Value",
    [
        "Lower chargeback losses via earlier detection",
        "Faster response through risk-based prioritization",
        "Clear accountability across customer, merchant, and admin teams",
        "Improved audit readiness and policy consistency",
        "Scalable process for growing dispute volume",
    ],
    "Tie business value to measurable KPIs: loss rate, response SLA, and resolution rate.",
)

add_bullets_slide(
    "Next Steps",
    [
        "Add production reason-code policy mappings",
        "Integrate with case management and ticketing tools",
        "Automate daily run scheduling",
        "Track precision, false positives, and recovery rate",
        "Refine intervention thresholds by merchant segment",
    ],
    "Close with a phased 30/60/90 day rollout recommendation.",
)

prs.save(str(OUT_PATH))
print(str(OUT_PATH))
